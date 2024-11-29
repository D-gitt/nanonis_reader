from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
import numpy as np
import os
import glob
import nanonis_reader as nr
from mpl_toolkits.axes_grid1 import make_axes_locatable
import io

class NanonisData:
    def __init__(self, base_path, file_number=None, keyword=None):
        '''
        base_bath: filepath
        fine_number: file number (ex: 0015.sxm -> 15)
        keyword: file name filter (ex: Au111_xxx_0013.sxm -> 'Au')
        '''
        # 파일 확장자별 함수 매핑
        file_handlers = {
            '.sxm': nr.nanonis_sxm.Load,
            '.dat': nr.nanonis_dat.Load,
            '.3ds': nr.nanonis_3ds.Load
        }
        
        # 파일 번호가 주어진 경우
        if isinstance(file_number, (int, str)):
            # 숫자를 4자리 문자열로 변환 (예: 16 -> '0016')
            number_str = str(file_number).zfill(4)
            
            # 키워드가 주어진 경우
            if keyword:
                pattern = os.path.join(base_path, f'*{keyword}*_{number_str}.*')
            else:
                pattern = os.path.join(base_path, f'*_{number_str}.*')
                
            matching_files = glob.glob(pattern)
            
            if not matching_files:
                if keyword:
                    raise ValueError(f"No file found with number {number_str} and keyword '{keyword}'")
                else:
                    raise ValueError(f"No file found with number {number_str}")
            
            if len(matching_files) > 1:
                print(f"Warning: Multiple files found. Files found:")
                for f in matching_files:
                    print(f"- {os.path.basename(f)}")
                print(f"Using the first one: {os.path.basename(matching_files[0])}")
            
            filepath = matching_files[0]
        else:
            filepath = base_path
            
        # 파일 확장자 추출
        _, extension = os.path.splitext(filepath)
        
        # 해당하는 함수 찾아서 실행
        if extension in file_handlers:
            data = file_handlers[extension](filepath)
            # data의 모든 속성을 현재 객체(self)에 복사
            for attr_name in dir(data):
                if not attr_name.startswith('_'):  # private 속성은 제외
                    setattr(self, attr_name, getattr(data, attr_name))
        else:
            raise ValueError(f"Unsupported file extension: {extension}")

class DataToPPT:
    def __init__(self, base_path, start_num, end_num, keyword=None, output_filename='output.pptx'):
        '''
        base_path: 파일들이 있는 경로
        start_num: 시작 파일 번호
        end_num: 끝 파일 번호
        keyword: 파일 이름 필터 (예: 'Au', 'SiCG')
        output_filename: 생성될 PPT 파일 이름
        '''
        self.base_path = base_path
        self.start_num = start_num
        self.end_num = end_num
        self.keyword = keyword
        self.output_filename = output_filename
        self.prs = Presentation()

    def get_scan_parameters(self, data):
        '''
        파일 타입에 따라 적절한 파라미터 추출 함수 호출
        '''
        if data.fname.endswith('.sxm'):
            return self.get_sxm_parameters(data)
        elif data.fname.endswith('.dat'):
            return self.get_dat_parameters(data)
        elif data.fname.endswith('.3ds'):
            return self.get_3ds_parameters(data)
        else:
            raise ValueError(f"Unsupported file type: {data.fname}")

    def get_sxm_parameters(self, data):
        '''
        header에서 자주 사용되는 파라미터들을 추출하는 함수
        반환값: pixels, scan_range, scan_dir, bias, current 등
        '''
        def format_date(date_str):
            '''
            일.월.년 형식을 년.월.일 형식으로 변환
            '''
            try:
                day, month, year = date_str.split('.')
                return f"{year}.{month}.{day}"
            except:
                return date_str  # 파싱 실패시 원본 반환
            
        params = {
            'pixels': data.header['scan_pixels'],
            'range': data.header['scan_range'],
            'direction': data.header['scan_dir'],
            'angle': data.header['scan_angle'],
            'bias': data.header['bias>bias (v)'],
            'current': data.header['z-controller>setpoint'],
            'scan_time': data.header['rec_time'],
            'scan_date': format_date(data.header['rec_date']),
        }
        params['aspect_ratio'] = (params['pixels'][0]/params['pixels'][1])*(params['range'][1]/params['range'][0])
        
        return params
        
    def get_dat_parameters(self, data):
        '''
        .dat 파일의 파라미터 추출
        Z rel (m) 포함 여부에 따라 다른 파라미터 반환
        '''
        def format_date(date_str):
            '''
            일.월.년 시:분:초 형식을 년.월.일 시:분:초 형식으로 변환
            '''
            try:
                # 날짜와 시간 분리
                date_part, time_part = date_str.split(' ')
                
                # 날짜 부분 변환
                day, month, year = date_part.split('.')
                formatted_date = f"{year}.{month}.{day}"
                
                # 날짜와 시간 다시 합치기
                return f"{formatted_date}_{time_part}"
            except:
                return date_str  # 파싱 실패시 원본 반환

        if 'Z rel (m)' in data.signals.keys():
            params = {
                'bias': data.header['Bias>Bias (V)'],
                'current': data.header['Z-Controller>Setpoint'],
                'sweep_num': data.header['Bias Spectroscopy>Number of sweeps'],
                'offset': data.header['Z Spectroscopy>Initial Z-offset (m)'],
                'sweep_z': data.header['Z Spectroscopy>Sweep distance (m)'],
                'sweep_num': data.header['Z Spectroscopy>Number of sweeps'],
                'comment': data.header['Comment01'],
                'saved_date': format_date(data.header['Saved Date']),
            }
            return params
        else:
            params = {
                'bias': data.header['Bias>Bias (V)'],
                'current': data.header['Z-Controller>Setpoint'],
                'sweep_start': data.header['Bias Spectroscopy>Sweep Start (V)'],
                'sweep_end': data.header['Bias Spectroscopy>Sweep End (V)'],
                'sweep_num': data.header['Bias Spectroscopy>Number of sweeps'],
                'comment': data.header['Comment01'],
                'saved_date': format_date(data.header['Saved Date']),
            }
            return params
        
        

    def get_3ds_parameters(self, data):
        '''
        .3ds 파일의 파라미터 추출
        '''
        params = {}
        # .3ds 파일에 필요한 파라미터들
        return params
    
    def get_sxm_info_text(self, params):
        '''
        .sxm 파일의 정보 텍스트 생성
        '''
        info_texts = []
        info_texts.append(f"{float(params['bias'])} V /")
        current = float(params['current'])
        if abs(current) >= 1e-9:
            # nA 단위로 표시
            info_texts.append(f"{current*1e9:.0f} nA")
        else:
            # pA 단위로 표시
            info_texts.append(f"{current*1e12:.0f} pA")
        info_texts.append(f"\n{params['range'][0]*1e9:.0f} x {params['range'][1]*1e9:.0f} nm²")
        info_texts.append(f"({params['direction']}, {float(params['angle']):.1f}˚)")
        info_texts.append(f"\n({params['scan_date']}_{params['scan_time']})")
        
        return " ".join(info_texts)

    def get_dat_info_text(self, params):
        '''
        .dat 파일의 정보 텍스트 생성
        '''
        if 'sweep_z' in params:
            info_texts = []
            info_texts.append(f"{float(params['bias'])} V /")
            current = float(params['current'])
            if abs(current) >= 1e-9:
                # nA 단위로 표시
                info_texts.append(f"{current*1e9:.0f} nA")
            else:
                # pA 단위로 표시
                info_texts.append(f"{current*1e12:.0f} pA")
            info_texts.append(f"\n({params['saved_date']})")
        else:
            info_texts = []
            info_texts.append(f"{float(params['bias'])} V /")
            current = float(params['current'])
            if abs(current) >= 1e-9:
                # nA 단위로 표시
                info_texts.append(f"{current*1e9:.0f} nA")
            else:
                # pA 단위로 표시
                info_texts.append(f"{current*1e12:.0f} pA")
            info_texts.append(f"\n{float(params['sweep_start'])} V to {float(params['sweep_end'])} V (sweeps: {params['sweep_num']})")
            info_texts.append(f"\n({params['saved_date']})")
        
        return " ".join(info_texts)

    def get_3ds_info_text(self, params):
        '''
        .3ds 파일의 정보 텍스트 생성
        '''
        # .3ds 파일에 맞는 정보 포맷
        return "3DS file parameters"

    def get_3sigma_limits(self, data):
        mean = np.nanmean(data)
        sigma = np.nanstd(data)
        return mean + np.array([-3, 3]) * sigma

    def process_sxm_file(self, data):
        '''
        .sxm 파일 처리 함수
        '''
        params = self.get_scan_parameters(data)

        base_size = 5
        figsize = (base_size, base_size)

        # 첫 번째 이미지 (topography)
        fig = plt.figure(figsize=figsize)
        ax = fig.add_subplot(111)      
        
        topo = nr.nanonis_sxm.topography(data)
        z_data = topo.get_z('subtract linear fit', 'fwd')

        origin = 'upper' if params['direction'] == 'down' else 'lower'
        vmin, vmax = self.get_3sigma_limits(z_data)
        nanox = nr.cmap_custom.nanox()
        bwr = nr.cmap_custom.bwr()
        
        # 이미지 플롯
        im = ax.imshow(z_data, origin=origin, vmin=vmin, vmax=vmax, 
                    aspect=params['aspect_ratio'], cmap=nanox, interpolation='none')

        # colorbar 추가
        plt.draw()
        posn = ax.get_position()
        cax = fig.add_axes([posn.x1 + 0.01, posn.y0, 
                       0.02, posn.height])
        plt.colorbar(im, cax=cax)

        # figure를 이미지로 저장
        img_stream1 = io.BytesIO()
        plt.savefig(img_stream1, format='png', bbox_inches='tight', pad_inches=0.01)
        img_stream1.seek(0)
        plt.close()

        # 두 번째 이미지 (differentiated)
        fig = plt.figure(figsize=figsize)
        ax = fig.add_subplot(111)
        
        z_data_diff = topo.get_z('differentiate', 'fwd')
        vmin, vmax = self.get_3sigma_limits(z_data_diff)
        im = ax.imshow(z_data_diff, origin=origin, vmin=vmin, vmax=vmax, 
                    aspect=params['aspect_ratio'], cmap=nanox, interpolation='none')

        # plt.draw()
        # posn = ax.get_position()
        # cax = fig.add_axes([posn.x1 + 0.01, posn.y0, 0.02, posn.height])
        # plt.colorbar(im, cax=cax)

        plt.draw()
        posn = ax.get_position()
        cax = fig.add_axes([posn.x1 + 0.01, posn.y0, 0.02, posn.height])
        cbar = plt.colorbar(im, cax=cax)
        cbar.formatter.set_powerlimits((-3, 4))  # scientific notation 사용 범위 설정
        cbar.update_ticks()

        img_stream2 = io.BytesIO()
        plt.savefig(img_stream2, format='png', bbox_inches='tight', pad_inches=0.01)
        img_stream2.seek(0)
        plt.close()

        if 'LI_Demod_1_X' in data.signals.keys():
            fig = plt.figure(figsize=figsize)
            ax = fig.add_subplot(111)

            didv = nr.nanonis_sxm.didvmap(data)
            didv_data = didv.get_map()
            vmin, vmax = self.get_3sigma_limits(didv_data)
            im = ax.imshow(didv_data, origin=origin, vmin=vmin, vmax=vmax, 
                        aspect=params['aspect_ratio'], cmap=bwr, interpolation='none')

            plt.draw()
            posn = ax.get_position()
            cax = fig.add_axes([posn.x1 + 0.01, posn.y0, 0.02, posn.height])
            plt.colorbar(im, cax=cax)

            img_stream3 = io.BytesIO()
            plt.savefig(img_stream3, format='png', bbox_inches='tight', pad_inches=0.01)
            img_stream3.seek(0)
            plt.close()
            
            return img_stream1, img_stream2, img_stream3

        return img_stream1, img_stream2
    
    def process_dat_file(self, data):
        '''
        .dat 파일 처리 함수
        '''
        params = self.get_dat_parameters(data)
        base_size = 5
        figsize = (base_size, base_size)

        if 'sweep_z' in params:
            spec = nr.nanonis_dat.z_spectrum(data)

            # Z-I linear
            spec_z = spec.get_iz()
            plt.figure(figsize=figsize)
            plt.plot(spec_z[0] * 1e9, spec_z[1] * 1e9, 'k-')
            plt.xlabel('Z (nm)')
            plt.ylabel('Current (nA)')
            img_stream1 = io.BytesIO()
            plt.savefig(img_stream1, format='png', bbox_inches='tight', pad_inches=0.1)
            img_stream1.seek(0)
            plt.close

            # Z-I log
            plt.figure(figsize=figsize)
            plt.plot(spec_z[0] * 1e9, np.abs(spec_z[1] * 1e9), 'k-')
            plt.xlabel('Z (nm)')
            plt.ylabel('|Current| (nA)')
            plt.yscale('log')
            plt.grid(True)
            img_stream2 = io.BytesIO()
            plt.savefig(img_stream2, format='png', bbox_inches='tight', pad_inches=0.1)
            img_stream2.seek(0)
            plt.close

            return img_stream1, img_stream2
        
        else:
            spec = nr.nanonis_dat.spectrum(data)

            # Scaled dI/dV
            didv_scaled = spec.didv_scaled()
            plt.figure(figsize=figsize)
            plt.plot(didv_scaled[0], didv_scaled[1] * 1e9, 'k-')
            plt.xlabel('Bias (V)')
            plt.ylabel('dI/dV (nS)')
            img_stream1 = io.BytesIO()
            plt.savefig(img_stream1, format='png', bbox_inches='tight', pad_inches=0.1)
            img_stream1.seek(0)
            plt.close

            # Normalized dI/dV
            didv_norm = spec.didv_normalized()
            plt.figure(figsize=figsize)
            plt.plot(didv_norm[0], didv_norm[1], 'k-')
            plt.xlabel('Bias (V)')
            plt.ylabel('Norm. dI/dV')
            img_stream2 = io.BytesIO()
            plt.savefig(img_stream2, format='png', bbox_inches='tight', pad_inches=0.1)
            img_stream2.seek(0)
            plt.close

            # I-V
            iv = spec.iv_raw()
            plt.figure(figsize=figsize)
            plt.plot(iv[0], iv[1] * 1e12, 'k-')
            plt.xlabel('Bias (V)')
            plt.ylabel('Current (pA)')
            img_stream3 = io.BytesIO()
            plt.savefig(img_stream3, format='png', bbox_inches='tight', pad_inches=0.1)
            img_stream3.seek(0)
            plt.close


            return img_stream1, img_stream2, img_stream3    
    
    def process_3ds_file(self, data):
        '''
        .3ds 파일 처리 함수
        '''
        # 3D 데이터 처리
        plt.figure(figsize=(10, 8))
        # 여기에 3D 데이터 처리 코드 추가
        
        img_stream = io.BytesIO()
        plt.savefig(img_stream, format='png', bbox_inches='tight')
        img_stream.seek(0)
        plt.close()
        return img_stream
    
    def add_slide(self, data):
        '''
        데이터를 처리하고 슬라이드에 추가하는 함수
        '''
        # 새 슬라이드 추가
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])  # 빈 슬라이드
        
        # 제목 추가
        title_shape = slide.shapes.title
        title_shape.text = f"File: {data.fname}"
        
        # 파일 확장자에 따른 처리
        if data.fname.endswith('.sxm'):
            params = self.get_scan_parameters(data)
            
            base_size = 3.2  # 이미지 크기를 조금 줄여서 3개가 들어갈 수 있게 조정
            width = Inches(base_size)
            height = Inches(base_size)

            img_top = Inches(1.5)

            if 'LI_Demod_1_X' in data.signals.keys():
                img_stream1, img_stream2, img_stream3 = self.process_sxm_file(data)
                # 세 이미지를 나란히 배치
                slide.shapes.add_picture(img_stream1, Inches(0), img_top,
                                    width=width)
                slide.shapes.add_picture(img_stream2, Inches(0 + base_size * 1.02), img_top,
                                    width=width)
                slide.shapes.add_picture(img_stream3, Inches(0 + base_size * 2.04), img_top,
                                    width=width)
            else:
                img_stream1, img_stream2 = self.process_sxm_file(data)
                # 두 이미지를 나란히 배치
                slide.shapes.add_picture(img_stream1, Inches(0), img_top,
                                    width=width)
                slide.shapes.add_picture(img_stream2, Inches(0 + base_size * 1.02), img_top,
                                    width=width)
            
            text_top = img_top + height + Inches(0.2)  # 0.2인치 간격
            info_text = self.get_sxm_info_text(params)

        elif data.fname.endswith('.dat'):
            params = self.get_dat_parameters(data)
            
            base_size = 3.2
            width = Inches(base_size)
            height = Inches(base_size)
            img_top = Inches(1.5)

            if 'sweep_z' in params:
                img_stream1, img_stream2 = self.process_dat_file(data)
                slide.shapes.add_picture(img_stream1, Inches(0), img_top, width=width)
                slide.shapes.add_picture(img_stream2, Inches(0 + base_size * 1.02), img_top, width=width)
            else:
                img_stream1, img_stream2, img_stream3 = self.process_dat_file(data)
                slide.shapes.add_picture(img_stream1, Inches(0), img_top, width=width)
                slide.shapes.add_picture(img_stream2, Inches(0 + base_size * 1.02), 
                                         img_top, width=width)
                slide.shapes.add_picture(img_stream3, Inches(0 + base_size * 2.04), 
                                         img_top, width=width)

            text_top = img_top + height + Inches(0.2)  # 0.2인치 간격
            info_text = self.get_dat_info_text(params)
            
        elif data.fname.endswith('.3ds'):
            params = self.get_scan_parameters(data)
            info_text = self.get_3ds_info_text(params)
            text_top = Inches(6)
            img_stream = self.process_3ds_file(data)
        
        # 추가 정보 텍스트 박스
        # left, top, width, height
        txBox = slide.shapes.add_textbox(Inches(1), text_top,
                                    Inches(8), Inches(0.5))
        tf = txBox.text_frame
        tf.text = info_text if info_text else "No parameters available"

    def generate_ppt(self):
        '''
        PPT 생성 메인 함수
        '''
        print(f"Generating PPT for files {self.start_num} to {self.end_num}...")
        
        for i in range(self.start_num, self.end_num + 1):
            try:
                # 파일 로드
                data = NanonisData(self.base_path, i, self.keyword)
                print(f"Processing file: {data.fname}")
                
                # 슬라이드 추가
                self.add_slide(data)
                
            except ValueError as e:
                print(f"Skipping number {i}: {str(e)}")
                continue
        
        # PPT 저장
        save_path = self.base_path + 'PPT/'
        self.prs.save(save_path + self.output_filename)
        print(f"PPT has been saved as {save_path + self.output_filename}")